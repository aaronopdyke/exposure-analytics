"""Build the WBG-branded comparison deck (Morocco) with python-pptx.

Starts from the official WBG template; keeps its Cover and the light
acknowledgements/disclaimer slide (wording lives on the layout, preserved
verbatim), deletes the other demo slides, and builds clean content slides:
white canvas + manually placed title/images/text (no body placeholders, no
split layouts - per formatting guidance for this deck). Seven models, all
monetary figures in constant latest-year US$; comparison numbers computed
from docs/data/MAR/meta.json at build time.

Usage: py site/tools/build_deck.py [out.pptx]
"""

import json
import os
import shutil
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SITE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(SITE, 'figures')
DATA = os.path.join(os.path.dirname(SITE), 'docs', 'data', 'MAR')
TEMPLATE = os.path.expanduser(
    r'~\.claude\skills\wbg-deck-template\assets\WBG_Template.pptx')
OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
    SITE, 'Morocco_Exposure_Model_Comparison.pptx')

NAVY = RGBColor(0x17, 0x40, 0x6D)
BLUE = RGBColor(0x0F, 0x6F, 0xC6)
GREY = RGBColor(0x40, 0x40, 0x40)
BRIGHT = RGBColor(0x00, 0x9D, 0xD9)
CYAN = RGBColor(0x0B, 0xD0, 0xD9)
GREEN = RGBColor(0x10, 0xCF, 0x9B)
OLIVE = RGBColor(0xA5, 0xC2, 0x49)
TINT = RGBColor(0xDB, 0xEF, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

meta = json.load(open(os.path.join(DATA, 'meta.json'), encoding='utf-8'))
T = meta['totals']
YR = meta['deflator']['target_year']
CTY = meta['country'].split(' (')[0]
# vintages/growth used by the vintage-reconciliation slide
USD = f'constant {YR} US$'
B = lambda v: f'${v / 1e9:,.0f} billion'
MM = lambda v: f'{v / 1e6:,.0f} Mm²'
ML = lambda v: f'{v / 1e6:,.1f} million'
pct = lambda a, b: f'{100 * a / (a + b):,.0f}%'


def keep_only(prs, keep_idx):
    ids = list(prs.slides._sldIdLst)
    for i, sld in enumerate(ids):
        if i not in keep_idx:
            rid = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            prs.part.drop_rel(rid)
            prs.slides._sldIdLst.remove(sld)


def layout_by_name(prs, name):
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name.strip() == name:
                return lay
    raise LookupError(name)


def text_block(slide, x, y, w, lines, h=None):
    """lines: list of (text, size, bold, color, space_after_pt)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w),
                                  Inches(h or 1))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, size, bold, color, gap) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = txt
        r.font.name = 'Arial'
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def content_slide(prs, layout, title):
    """Clean white content slide: the template's content layouts are all
    half-split art panels (explicitly not wanted for this deck), so every
    placeholder is dropped, the layout art is covered with a white canvas,
    and the title + rule are drawn manually in brand type."""
    from pptx.enum.shapes import MSO_SHAPE
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    canvas = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    canvas.line.fill.background()
    canvas.shadow.inherit = False
    text_block(s, 0.45, 0.32, 12.4, [(title, 24, True, NAVY, 0)])
    rule = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.02), Inches(12.35),
        Pt(1.6))
    rule.fill.solid()
    rule.fill.fore_color.rgb = NAVY
    rule.line.fill.background()
    rule.shadow.inherit = False
    return s


def pic_fit(slide, path, x, y, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    return slide.shapes.add_picture(
        path, Inches(x), Inches(y),
        width=Emu(int(Inches(1) * iw * scale)),
        height=Emu(int(Inches(1) * ih * scale)))


def model_slide(prs, layout, title, img, headline, sub, bullets):
    s = content_slide(prs, layout, title)
    pic_fit(s, os.path.join(FIG, img), 0.35, 1.35, 6.9, 5.7)
    lines = [(headline, 30, True, NAVY, 2), (sub, 13, False, BLUE, 16)]
    for b in bullets:
        lines.append(('–  ' + b, 12.5, False, GREY, 8))
    text_block(s, 7.5, 1.6, 5.5, lines)
    return s


def section_slide(prs, layout, num, title):
    """Navy full-bleed section divider."""
    from pptx.enum.shapes import MSO_SHAPE
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    canvas = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = NAVY
    canvas.line.fill.background()
    canvas.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(2.85),
                             Inches(1.5), Pt(3.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    bar.shadow.inherit = False
    text_block(s, 0.9, 3.05, 11.6, [
        (f'SECTION {num}', 13, True, CYAN, 8),
        (title, 38, True, WHITE, 0)])
    return s


shutil.copy2(TEMPLATE, OUT)
prs = Presentation(OUT)
# keep Cover (idx 0) and the light full-text disclaimer variant (idx 21,
# layout "Acknowledgements and disclaimers 7" - wording lives on the layout
# and is preserved untouched)
keep_only(prs, {0, 21})
lay = layout_by_name(prs, 'Content slide 1')

# ---- cover ----------------------------------------------------------------
# match placeholders by idx: python-pptx yields fresh proxies on every
# iteration, so identity checks silently match nothing
for ph in prs.slides[0].placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text_frame.text = 'Comparing Building Exposure Models'
    elif ph.has_text_frame:
        ph.text_frame.text = CTY

# ---- section 1: the comparison --------------------------------------------
section_slide(prs, lay, 1, 'Comparing the totals')

native = [T['gar15_usd_adj'], T['giri_total_usd_adj'],
          T['gem23_bldg_repl_usd_adj'], T['gem26_bldg_repl_usd_adj']]
ratio = max(native) / min(native)
s = content_slide(prs, lay, 'Country totals by model')
pic_fit(s, os.path.join(FIG, 'bars_totals.png'), 0.6, 1.45, 12.1, 4.35)
text_block(s, 0.75, 5.95, 11.8, [
    (f"After price-year alignment: GAR15 {B(T['gar15_usd_adj'])} · GIRI "
     f"{B(T['giri_total_usd_adj'])} · GEM v2023 "
     f"{B(T['gem23_bldg_repl_usd_adj'])} · GEM v2026 "
     f"{B(T['gem26_bldg_repl_usd_adj'])}. All figures are structures-only "
     '(GEM contents dropped). The remaining spread is method, not price '
     'levels - see the vintage slide.', 12.5, False, GREY, 8),
    ('With GBA-calibrated storeys the three footprint estimates now sit '
     f"between {B(min(T['overture_value_usd_adj'], T['msb_value_usd_adj']))} "
     f"and {B(T['gba_value_usd_adj'])} - same buildings, different height "
     'treatment.', 12.5, False, GREY, 0),
])

VG = meta['vintages']['value_growth']
FG = meta['vintages']['floor_growth']
s = content_slide(prs, lay,
                  f'Country totals, grown to {YR} for growth since each '
                  'model\'s vintage')
pic_fit(s, os.path.join(FIG, 'bars_totals_grown.png'), 0.6, 1.45, 12.1, 4.35)
text_block(s, 0.75, 5.95, 11.8, [
    ('Identical charts to the previous slide, but each model is grown from '
     'its data vintage: values along the WB produced-capital path (CWON '
     f"NW.PCA.TO, ~{VG['trailing_cagr']:.1%}/yr recently), floor areas at "
     f"the GHSL built-up-volume rate (~{FG['rate']:.1%}/yr). "
     f"GAR15 (2011 inputs, x{VG['factors']['gar15']:.2f}) jumps to "
     f"{B(T['gar15_usd_adj'] * VG['factors']['gar15'])}.",
     12.5, False, GREY, 8),
    ('Growing every model to a common stock year makes the spread WIDER, '
     'not narrower - the disagreements are about method, not vintage. '
     'Factors of x1.00 mean the data already reflects ' + str(YR) + '.',
     12.5, False, GREY, 0),
])

s = content_slide(prs, lay, 'Regional breakdown by WB Admin-1 unit')
pic_fit(s, os.path.join(FIG, 'adm1_grouped.png'), 0.5, 1.4, 7.6, 5.6)
text_block(s, 8.35, 1.7, 4.6, [
    ('Same top regions in all three models', 15, True, NAVY, 8),
    ('–  Casablanca–Settat leads in every model, followed by '
     'Rabat–Salé–Kénitra and Marrakech–Safi.', 12.5, False, GREY, 7),
    ('–  Disagreement concentrates in secondary regions — GEM '
     'is notably lower in the south and east.', 12.5, False, GREY, 7),
    ('–  Western Sahara (hatched unit): small in all models; GEM '
     'values come from its two southern regions via the boundary '
     'crosswalk.', 12.5, False, GREY, 7),
])

# ---- section 2: the models ------------------------------------------------
section_slide(prs, lay, 2, 'The models evaluated')

# ---- overview: what is being compared -------------------------------------
def _box(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                               Inches(y), Inches(w), Inches(h))
    b.adjustments[0] = 0.08
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(1)
    b.shadow.inherit = False
    return b


def _card(slide, x, y, w, name, color, desc, h=1.28):
    _box(slide, x, y, w, h, WHITE, line=RGBColor(0xD5, 0xDE, 0xE6))
    strip = _box(slide, x, y, 0.09, h, color)
    strip.adjustments[0] = 0.5
    text_block(slide, x + 0.22, y + 0.12, w - 0.35, [
        (name, 13.5, True, NAVY, 2),
        (desc, 10.5, False, GREY, 0),
    ])


s = content_slide(prs, lay, 'Datasets evaluated')
_COLS = [
    ('CAPITAL STOCK OF BUILDINGS', 'top-down, from WB produced capital', [
        ('GAR15 — Global Assessment Report 2015', BLUE,
         'UN Int. Strategy for Disaster Reduction (UNISDR); ~5 km point '
         'grid, ~2011 inputs; native urban/rural split'),
        ('GIRI — Global Infrastructure Resilience Index', NAVY,
         'Building Exposure Model, UNEP/GRID-Geneva (~2023); 5 km rasters, '
         '~2020 inputs; res / non-res layers'),
    ]),
    ('BUILDING REPLACEMENT COST', 'bottom-up engineering, excl. contents', [
        ('GEM v2023.1.1 — Global Exposure Model', BRIGHT,
         'Global Earthquake Model Foundation; ADM1 x occupancy summaries; '
         '2021 US$ price base'),
        ('GEM v2026.0.0 — Global Exposure Model', BRIGHT,
         'stock updated to 2025, costs re-based to 2024-25'),
    ]),
    ('ESTIMATED REPLACEMENT VALUE', 'floor area x GEM unit cost (UCC)', [
        ('Overture Maps (2026-06)', CYAN,
         'mapped footprints: OpenStreetMap + Microsoft + Google'),
        ('GBA — Global Building Atlas', GREEN,
         'Technical University of Munich; LoD1 footprints WITH modelled '
         'heights (~2024)'),
        ('Microsoft GlobalML (2026-02)', OLIVE,
         'machine-learned footprints; partial coverage in Morocco'),
    ]),
]
_x0, _w, _gap = 0.5, 4.0, 0.18
for _i, (_hdr, _sub, _cards) in enumerate(_COLS):
    _x = _x0 + _i * (_w + _gap)
    _box(s, _x, 1.35, _w, 0.78, NAVY)
    text_block(s, _x + 0.18, 1.44, _w - 0.3, [
        (_hdr, 13, True, WHITE, 1), (_sub, 10, False, TINT, 0)])
    for _j, (_name, _color, _desc) in enumerate(_cards):
        _card(s, _x, 2.32 + _j * 1.44, _w, _name, _color, _desc)
_box(s, _x0, 6.85, 3 * _w + 2 * _gap, 0.52, TINT)
text_block(s, _x0 + 0.2, 6.95, 3 * _w + 2 * _gap - 0.4, [
    (f'All seven aligned to WB Official Boundaries Admin-1 · {USD} · '
     'structures only (contents excluded everywhere)', 12, True, NAVY, 0)])

# ---- one slide per model --------------------------------------------------
model_slide(prs, lay, 'GAR15 — UNISDR Global Exposure Dataset (2015)',
    'map_gar15.png', B(T['gar15_usd_adj']),
    f'capital stock of buildings, {USD}', [
    f"{B(T['gar15_usd'])} in its native 2005 US$ (WB produced capital, "
    'structures - no contents)',
    f"urban {B(T['gar15_urban_usd_adj'])} · rural "
    f"{B(T['gar15_rural_usd_adj'])} "
    f"({pct(T['gar15_urban_usd'], T['gar15_rural_usd'])} urban, native split)",
    'top-down: national capital stock distributed on a ~5 km grid '
    '(1 km at the coast) using population and GDP proxies',
    'source: HDX per-country shapefiles; Western Sahara is a separate '
    'GAR15 dataset, shown as the hatched WB unit',
])
model_slide(prs, lay, 'GIRI — Building Exposure Model (~2023)',
    'map_giri.png', B(T['giri_total_usd_adj']),
    f'capital stock of buildings, {USD}', [
    f"{B(T['giri_total_usd'])} in its native 2018 US$ (WB produced "
    'capital, structures - no contents)',
    f"residential {B(T['giri_res_usd_adj'])} · non-residential "
    f"{B(T['giri_nres_usd_adj'])}",
    'top-down 5×5 km global rasters by UNEP/GRID-Geneva; includes '
    'social infrastructure (health, education)',
    'no urban/rural split published; res / non-res layers instead',
])
model_slide(prs, lay, 'GEM — Global Exposure Model (v2023.1.1 and v2026.0.0)',
    'map_gem26.png', B(T['gem26_bldg_repl_usd_adj']),
    f'v2026 building replacement cost, {USD} (excl. contents)', [
    f"v2023.1.1: {B(T['gem23_bldg_repl_usd_adj'])} in {USD} "
    '(2021 US$ native; no urban/rural split published)',
    f"v2026.0.0: costs re-based to 2024-25, stock updated to 2025; "
    f"{ML(T['gem26_buildings'])} buildings · "
    f"{MM(T['gem26_floor_area_m2'])} floor area",
    'most of the v2023 -> v2026 increase is consistent with ~3 years of '
    'stock growth plus the cost re-basing (see vintage slide)',
    f"GEM's own published total is {B(T['gem26_bldg_repl_usd'] + T['gem26_contents_usd'])} "
    f"INCLUDING {B(T['gem26_contents_usd'])} contents - dropped here so all "
    'models are structures-only',
    "bottom-up engineering model; GEM's two southern regions map to the "
    'WB Western Sahara unit via region matching (map shows v2026)',
])
model_slide(prs, lay, 'Overture Maps — building footprints (2026-06)',
    'map_overture.png', B(T['overture_value_usd_adj']),
    f'ESTIMATED replacement value = floor area × GEM UCC, {USD}', [
    f"{ML(T['overture_buildings'])} mapped buildings · "
    f"{MM(T['overture_floor_area_m2'])} estimated floor area",
    'num_floors/height rarely populated in Morocco: reported attributes '
    'used where present, GBA-calibrated storey multipliers (per ADM1 x '
    'urban/rural) for the rest',
    'community/corporate footprints (OSM, Microsoft, Google) from the '
    'pinned S3 release, aggregated in DuckDB',
    f"urban share {pct(T['overture_floor_area_m2_urban'], T['overture_floor_area_m2_rural'])} "
    'of floor area (GHSL DEGURBA 2020)',
])
model_slide(prs, lay, 'Global Building Atlas — LoD1 (~2024)',
    'map_gba.png', B(T['gba_value_usd_adj']),
    f'ESTIMATED replacement value = floor area × GEM UCC, {USD}', [
    f"{ML(T['gba_buildings'])} buildings · "
    f"{MM(T['gba_floor_area_m2'])} floor area · "
    f"{T['gba_volume_m3'] / 1e9:,.1f} km³ volume",
    'the only footprint source with modelled heights (TUM LoD1) — '
    'hence the largest floor area and estimated replacement value',
    'floor area via 3-tier storey heights (3.0 / 3.25 / 3.5 m); '
    'missing heights counted as one storey',
    f"urban share {pct(T['gba_floor_area_m2_urban'], T['gba_floor_area_m2_rural'])} "
    'of floor area (GHSL DEGURBA 2020)',
])
model_slide(prs, lay, 'Microsoft Buildings — GlobalML footprints (2026-02)',
    'map_msb.png', B(T['msb_value_usd_adj']),
    f'ESTIMATED replacement value = floor area × GEM UCC, {USD}', [
    f"{ML(T['msb_buildings'])} buildings · "
    f"{MM(T['msb_footprint_m2'])} footprint · "
    f"{MM(T['msb_floor_area_m2'])} estimated floor area",
    'native Microsoft GlobalML release; its height field exists but is '
    'unmodelled (-1) throughout Morocco, so storeys come from the '
    'GBA calibration',
    'coverage is PARTIAL in Morocco: no tiles over Souss-Massa, Beni '
    'Mellal-Khenifra and most of Draa-Tafilalet - totals undercount',
    f"urban share {pct(T['msb_floor_area_m2_urban'], T['msb_floor_area_m2_rural'])} "
    'of floor area (GHSL DEGURBA 2020)',
])

# ---- section 3: assumptions & sensitivity ---------------------------------
section_slide(prs, lay, 3, 'Assumptions & sensitivity')

vin = meta['vintages']['years']
gar_grown = T['gar15_usd_adj'] * VG['factors']['gar15']
s = content_slide(prs, lay, 'Does data age explain the gaps? Mostly not')
pic_fit(s, os.path.join(FIG, 'vintage_growth.png'), 0.6, 1.45, 12.1, 4.35)
text_block(s, 0.75, 5.95, 11.8, [
    ('Growth rates are data-driven, not assumed: values follow the WB '
     'Wealth Accounts produced-capital series for Morocco (NW.PCA.TO - the '
     'same CWON data GAR15 and GIRI derive from), floor areas the GHSL '
     'built-up-volume rate. Growing GAR15 along its own source series to '
     f'{YR} yields {B(gar_grown)} - OVERSHOOTING GIRI '
     f"({B(T['giri_total_usd_adj'] * VG['factors']['giri'])}) rather than "
     'converging: their building-share methods differ, not their vintage.',
     12.5, False, GREY, 8),
    ('(Floor areas grow far more slowly - GHSL built-up volume, ~'
     f"{FG['rate']:.1%}/yr - so age barely moves them.) The exception "
     'remains GEM v2023 -> v2026, where growth plus cost re-basing '
     'explains most of the change.', 12.5, False, GREY, 0),
])

# ---- storey-rule sensitivity ---------------------------------------------
SENS = meta['sensitivity']


def _sens_range(k):
    b = SENS['baseline']['models'][k]['value_usd_adj']
    vals = [s['models'][k]['value_usd_adj'] for s in SENS.values()]
    return min(vals) / b - 1, max(vals) / b - 1


_rngs = [_sens_range(k) for k in ('overture', 'gba', 'msb')]
sw_lo = min(r[0] for r in _rngs)
sw_hi = max(r[1] for r in _rngs)
s = content_slide(prs, lay, 'Sensitivity analysis: storey height assumptions')
pic_fit(s, os.path.join(FIG, 'bars_sensitivity.png'), 0.6, 1.45, 12.1, 4.35)
text_block(s, 0.75, 5.95, 11.8, [
    (f'Extreme rules move footprint-model totals {sw_lo:+.0%} to '
     f'{sw_hi:+.0%} (vs a ~{ratio:.1f}× spread between models). Four '
     'alternative height-to-floor-area rules, recomputed per building '
     'for GBA and carried through the GBA-calibrated multipliers for '
     'Overture/Microsoft: integer storeys rounded DOWN (min 1), rounded UP, '
     'and the 3-tier storey height (3.0 / 3.25 / 3.5 m) shifted by '
     '±0.25 m.', 12.5, False, GREY, 8),
    ('Even the extreme rules neither reorder the models nor close the gap '
     'to the produced-capital models - the comparison\'s conclusions are '
     'insensitive to the storey convention.', 12.5, False, GREY, 0),
])

s = content_slide(prs, lay, 'Urban vs rural: models broadly agree, definitions differ')
pic_fit(s, os.path.join(FIG, 'urban_rural.png'), 0.6, 1.5, 12.1, 4.4)
text_block(s, 0.75, 6.15, 11.8, [
    ('GAR15 and GEM v2026 publish their own urban/rural splits (GEM: residential '
     'only — COM/IND have no settlement tag). Overture, GBA and Microsoft are '
     'classified with GHSL DEGURBA 2020 (urban = suburban through '
     "urban-centre classes) at each building's 1 km grid cell.",
     12.5, False, GREY, 8),
    ('Urban shares cluster at 52–77% — Overture sits lowest because its '
     'missing height attributes flatten dense urban stock the most.',
     12.5, False, GREY, 0),
])

# ---- order: move the kept disclaimer slide to the end ---------------------
ids = list(prs.slides._sldIdLst)
prs.slides._sldIdLst.remove(ids[1])
prs.slides._sldIdLst.append(ids[1])

prs.save(OUT)
print('deck written:', OUT, f'({len(list(prs.slides._sldIdLst))} slides)')

# mirror to the Drive Outputs folder so the shared copy never goes stale
sys.path.insert(0, os.path.dirname(SITE))
try:
    from src.config import OUTPUTS_DIR
    if OUTPUTS_DIR and os.path.isdir(OUTPUTS_DIR):
        shutil.copy2(OUT, os.path.join(OUTPUTS_DIR, os.path.basename(OUT)))
        print('copied to', OUTPUTS_DIR)
except Exception as e:  # Drive offline is fine - the repo copy is canonical
    print('Drive copy skipped:', e)
